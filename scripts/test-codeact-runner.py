"""Smoke-test a published Bubblewrap Runner without Docker."""
import base64
import argparse
from contextlib import ExitStack
import io
import json
import os
from pathlib import Path
import secrets
import socket
import subprocess
import tempfile
import time
from urllib.error import HTTPError, URLError
from urllib.request import Request, urlopen
import zipfile


def request(url, payload=None, key=None):
    data = None if payload is None else json.dumps(payload).encode("utf-8")
    headers = {"Content-Type": "application/json"}
    if key:
        headers["Authorization"] = "Bearer " + key
    try:
        with urlopen(Request(url, data=data, headers=headers), timeout=150) as response:
            return response.status, response.read()
    except HTTPError as error:
        return error.code, error.read()


def unused_port():
    with socket.socket() as listener:
        listener.bind(("127.0.0.1", 0))
        return listener.getsockname()[1]


def main():
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument('--environment-file', type=Path, help='Test an already running service using its environment file.')
    arguments = parser.parse_args()
    runner = os.environ.get("CODEACT_RUNNER_DLL")
    python = os.environ.get("CODEACT_TEST_PYTHON", "/opt/openagent-code/venv/bin/python")
    if not arguments.environment_file and (not runner or not Path(runner).is_file()):
        raise SystemExit("Set CODEACT_RUNNER_DLL to a published OpenAgent.Runner.dll.")

    with ExitStack() as stack:
        log = stack.enter_context(tempfile.TemporaryFile())
        process = None
        if arguments.environment_file:
            environment = dict(line.split('=', 1) for line in arguments.environment_file.read_text().splitlines()
                               if line and not line.startswith('#'))
            base_url = environment['ASPNETCORE_URLS'].rstrip('/')
            key = environment['Runner__ApiKey']
            directory = environment['Runner__WorkspaceRoot']
        else:
            directory = stack.enter_context(tempfile.TemporaryDirectory(prefix="codeact-smoke-"))
            key = secrets.token_hex(32)
            base_url = f"http://127.0.0.1:{unused_port()}"
            environment = dict(
                os.environ,
                ASPNETCORE_URLS=base_url,
                Runner__ApiKey=key,
                Runner__WorkspaceRoot=directory,
                Runner__BubblewrapPath=os.environ.get("CODEACT_TEST_BWRAP", "/usr/bin/bwrap"),
                Runner__PythonPath=python,
            )
            process = subprocess.Popen(["dotnet", runner], env=environment, stdout=log, stderr=log)
        try:
            for _ in range(60):
                try:
                    status, _ = request(base_url + "/health")
                    if status == 200:
                        break
                except URLError:
                    pass
                if process is not None and process.poll() is not None:
                    raise AssertionError("Runner exited before becoming ready.")
                time.sleep(0.5)
            else:
                raise AssertionError("Runner did not become ready.")

            status, _ = request(base_url + "/v1/execute", {"code": "print(42)"})
            assert status == 401, status
            code = """
import os
import subprocess
from pathlib import Path
from openpyxl import Workbook
from pptx import Presentation
assert os.getuid() == 65532
assert 'Runner__ApiKey' not in os.environ
assert not os.path.exists('/var/run/docker.sock')
w = Workbook()
w.active['A1'] = 42
w.save('/output/report.xlsx')
p = Presentation()
p.slides.add_slide(p.slide_layouts[0]).shapes.title.text = 'Isolated CodeAct'
p.save('/output/report.pptx')
conversion = subprocess.run(['libreoffice', '--headless', '-env:UserInstallation=file:///tmp/lo', '--convert-to', 'pdf', '--outdir', '/output', '/output/report.pptx'], capture_output=True, timeout=90)
assert conversion.returncode == 0, conversion.stderr
assert Path('/output/report.pdf').read_bytes().startswith(b'%PDF')
print('isolated execution passed')
"""
            status, body = request(base_url + "/v1/execute", {"code": code}, key)
            assert status == 200, body.decode("utf-8", errors="replace")
            response = json.loads(body)
            assert response["exitCode"] == 0, response["stderr"]
            assert {file["name"] for file in response["files"]} == {"report.xlsx", "report.pptx", "report.pdf"}
            for artifact in response["files"]:
                if artifact['name'].endswith('.pdf'):
                    assert base64.b64decode(artifact['content']).startswith(b'%PDF')
                    continue
                with zipfile.ZipFile(io.BytesIO(base64.b64decode(artifact["content"]))) as archive:
                    assert archive.testzip() is None
            assert not list(Path(directory).iterdir()), "Task input directories were not removed."
            print("PASS: Runner authentication, Bubblewrap isolation, PPT/XLSX/PDF artifacts, and cleanup.")
        except BaseException:
            if process is not None:
                log.seek(0)
                print(log.read().decode("utf-8", errors="replace"))
            raise
        finally:
            if process is not None:
                process.terminate()
                try:
                    process.wait(timeout=15)
                except subprocess.TimeoutExpired:
                    process.kill()
                    process.wait(timeout=5)


if __name__ == "__main__":
    main()
